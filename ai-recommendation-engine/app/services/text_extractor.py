"""
Text Extractor Service
Extracts text content from uploaded files (PDF, DOCX, PPTX, TXT)
Used to automatically extract course content from uploaded documents.
"""

import logging
import io
from typing import Tuple

logger = logging.getLogger(__name__)


def extract_from_pdf(file_bytes: bytes) -> Tuple[str, int]:
    """Extract text from a PDF file. Returns (text, page_count)."""
    from PyPDF2 import PdfReader

    reader = PdfReader(io.BytesIO(file_bytes))
    pages = len(reader.pages)
    text_parts = []

    for page in reader.pages:
        page_text = page.extract_text()
        if page_text:
            text_parts.append(page_text.strip())

    return "\n\n".join(text_parts), pages


def extract_from_docx(file_bytes: bytes) -> Tuple[str, int]:
    """Extract text from a DOCX file. Returns (text, paragraph_count)."""
    from docx import Document

    doc = Document(io.BytesIO(file_bytes))
    text_parts = []

    for para in doc.paragraphs:
        if para.text.strip():
            text_parts.append(para.text.strip())

    # Also extract text from tables
    for table in doc.tables:
        for row in table.rows:
            row_text = " | ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
            if row_text:
                text_parts.append(row_text)

    return "\n\n".join(text_parts), len(text_parts)


def extract_from_pptx(file_bytes: bytes) -> Tuple[str, int]:
    """Extract text from a PPTX file. Returns (text, slide_count)."""
    from pptx import Presentation

    prs = Presentation(io.BytesIO(file_bytes))
    text_parts = []
    slide_count = len(prs.slides)

    for slide_num, slide in enumerate(prs.slides, 1):
        slide_texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    if paragraph.text.strip():
                        slide_texts.append(paragraph.text.strip())

        if slide_texts:
            text_parts.append(f"--- Slide {slide_num} ---\n" + "\n".join(slide_texts))

    return "\n\n".join(text_parts), slide_count


def extract_from_txt(file_bytes: bytes) -> Tuple[str, int]:
    """Extract text from a TXT file. Returns (text, line_count)."""
    # Try UTF-8 first, then latin-1 as fallback
    try:
        text = file_bytes.decode("utf-8")
    except UnicodeDecodeError:
        text = file_bytes.decode("latin-1")

    lines = len(text.strip().split("\n"))
    return text.strip(), lines


def extract_text(file_bytes: bytes, filename: str) -> dict:
    """
    Main extraction function. Auto-detects file type and extracts text.

    Returns:
        {
            "text": str,
            "filename": str,
            "pages": int,
            "file_type": str,
            "char_count": int
        }
    """
    filename_lower = filename.lower()

    try:
        if filename_lower.endswith(".pdf"):
            text, pages = extract_from_pdf(file_bytes)
            file_type = "PDF"
        elif filename_lower.endswith(".docx"):
            text, pages = extract_from_docx(file_bytes)
            file_type = "DOCX"
        elif filename_lower.endswith(".pptx"):
            text, pages = extract_from_pptx(file_bytes)
            file_type = "PPTX"
        elif filename_lower.endswith(".txt"):
            text, pages = extract_from_txt(file_bytes)
            file_type = "TXT"
        else:
            raise ValueError(f"Unsupported file type: {filename}")

        logger.info(
            f"Extracted {len(text)} chars from {filename} ({file_type}, {pages} pages/sections)"
        )

        return {
            "text": text,
            "filename": filename,
            "pages": pages,
            "file_type": file_type,
            "char_count": len(text),
        }

    except Exception as e:
        logger.error(f"Text extraction failed for {filename}: {str(e)}", exc_info=True)
        raise ValueError(f"Failed to extract text from {filename}: {str(e)}")
